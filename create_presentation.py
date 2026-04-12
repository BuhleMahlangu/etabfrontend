#!/usr/bin/env python3
"""
E-tab Learning Management System - Marketing Presentation Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
PRIMARY_BLUE = RGBColor(37, 99, 235)      # #2563EB
PRIMARY_PURPLE = RGBColor(124, 58, 237)   # #7C3AED
DARK_BG = RGBColor(15, 23, 42)            # #0F172A
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(248, 250, 252)      # #F8FAFC
ACCENT_GREEN = RGBColor(34, 197, 94)      # #22C55E
ACCENT_ORANGE = RGBColor(249, 115, 22)    # #F97316

def add_gradient_background(slide, color1, color2):
    """Add a gradient background to the slide"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = color1
    background.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_title_slide():
    """Title Slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    add_gradient_background(slide, DARK_BG, DARK_BG)
    
    # Logo placeholder area
    logo_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(1.5), Inches(2.3), Inches(2.3))
    logo_box.fill.solid()
    logo_box.fill.fore_color.rgb = PRIMARY_BLUE
    logo_box.line.fill.background()
    
    # Logo text
    logo_tf = logo_box.text_frame
    logo_tf.text = "E"
    logo_p = logo_tf.paragraphs[0]
    logo_p.font.size = Pt(72)
    logo_p.font.bold = True
    logo_p.font.color.rgb = WHITE
    logo_p.alignment = PP_ALIGN.CENTER
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
    tf = title_box.text_frame
    tf.text = "E-tab"
    p = tf.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11.333), Inches(0.8))
    tf2 = subtitle_box.text_frame
    tf2.text = "Next-Generation Learning Management System"
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(148, 163, 184)
    p2.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11.333), Inches(0.6))
    tf3 = tagline_box.text_frame
    tf3.text = "Empowering Education Through Technology"
    p3 = tf3.paragraphs[0]
    p3.font.size = Pt(18)
    p3.font.color.rgb = PRIMARY_BLUE
    p3.alignment = PP_ALIGN.CENTER

def add_overview_slide():
    """Project Overview Slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_BLUE
    header.line.fill.background()
    
    # Header title
    header_tf = header.text_frame
    header_tf.text = "What is E-tab?"
    hp = header_tf.paragraphs[0]
    hp.font.size = Pt(36)
    hp.font.bold = True
    hp.font.color.rgb = WHITE
    hp.alignment = PP_ALIGN.LEFT
    header_tf.margin_left = Inches(0.5)
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(2))
    tf = desc_box.text_frame
    tf.word_wrap = True
    tf.text = "E-tab is a comprehensive Learning Management System (LMS) designed for modern educational institutions. It connects learners, teachers, and administrators in a unified digital ecosystem."
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(71, 85, 105)
    p.line_spacing = 1.5
    
    # Key points
    points = [
        ("🏫", "Multi-School Support", "Manage multiple schools from a single platform"),
        ("👥", "Role-Based Access", "Separate dashboards for Learners, Teachers, and Admins"),
        ("📱", "Mobile-First Design", "Responsive interface that works on any device"),
        ("🤖", "AI-Powered Tutoring", "Built-in AI tutor to assist student learning"),
    ]
    
    y_pos = Inches(4)
    for icon, title, desc in points:
        # Icon box
        icon_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(0.6), Inches(0.6))
        icon_box.fill.solid()
        icon_box.fill.fore_color.rgb = PRIMARY_BLUE
        icon_box.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1.6), y_pos, Inches(4), Inches(0.4))
        ttf = title_box.text_frame
        ttf.text = title
        tp = ttf.paragraphs[0]
        tp.font.size = Pt(18)
        tp.font.bold = True
        tp.font.color.rgb = DARK_BG
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(1.6), y_pos + Inches(0.4), Inches(4), Inches(0.4))
        dtf = desc_box.text_frame
        dtf.text = desc
        dp = dtf.paragraphs[0]
        dp.font.size = Pt(12)
        dp.font.color.rgb = RGBColor(100, 116, 139)
        
        y_pos += Inches(0.9)

def add_tech_stack_slide():
    """Technology Stack Slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_PURPLE
    header.line.fill.background()
    
    header_tf = header.text_frame
    header_tf.text = "Technology Stack"
    hp = header_tf.paragraphs[0]
    hp.font.size = Pt(36)
    hp.font.bold = True
    hp.font.color.rgb = WHITE
    hp.alignment = PP_ALIGN.LEFT
    header_tf.margin_left = Inches(0.5)
    
    # Tech categories
    categories = [
        ("Frontend Framework", ["React 18", "React Router v6", "Vite 5"], PRIMARY_BLUE),
        ("Styling & UI", ["Tailwind CSS 3", "Custom CSS Animations", "Lucide Icons"], ACCENT_GREEN),
        ("State & Communication", ["React Context API", "Socket.IO Client", "Axios"], ACCENT_ORANGE),
        ("Development Tools", ["Vite Dev Server", "PostCSS", "Autoprefixer"], RGBColor(236, 72, 153)),
    ]
    
    x_pos = Inches(0.8)
    for cat_name, items, color in categories:
        # Category card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(2), Inches(2.8), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = color
        card.line.width = Pt(3)
        
        # Category header
        cat_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(2), Inches(2.8), Inches(0.8))
        cat_header.fill.solid()
        cat_header.fill.fore_color.rgb = color
        cat_header.line.fill.background()
        
        ch_tf = cat_header.text_frame
        ch_tf.text = cat_name
        chp = ch_tf.paragraphs[0]
        chp.font.size = Pt(14)
        chp.font.bold = True
        chp.font.color.rgb = WHITE
        chp.alignment = PP_ALIGN.CENTER
        
        # Items
        y_item = Inches(3.2)
        for item in items:
            item_box = slide.shapes.add_textbox(x_pos + Inches(0.2), y_item, Inches(2.4), Inches(0.5))
            itf = item_box.text_frame
            itf.text = "• " + item
            ip = itf.paragraphs[0]
            ip.font.size = Pt(14)
            ip.font.color.rgb = DARK_BG
            y_item += Inches(0.6)
        
        x_pos += Inches(3.1)

def add_features_slide():
    """Key Features Slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = ACCENT_GREEN
    header.line.fill.background()
    
    header_tf = header.text_frame
    header_tf.text = "Key Features"
    hp = header_tf.paragraphs[0]
    hp.font.size = Pt(36)
    hp.font.bold = True
    hp.font.color.rgb = WHITE
    hp.alignment = PP_ALIGN.LEFT
    header_tf.margin_left = Inches(0.5)
    
    features = [
        ("📚", "Learning Materials", "Upload and manage documents, videos, and resources with Cloudinary integration"),
        ("📝", "Quiz System", "Create timed quizzes with multiple question types, auto-grading, and anti-cheat measures"),
        ("✅", "Assignments", "Distribute, collect, and grade assignments with deadline tracking"),
        ("📊", "Progress Tracking", "Visual analytics and FET phase history for comprehensive student monitoring"),
        ("🔔", "Real-time Notifications", "Instant alerts via Socket.IO for announcements, deadlines, and messages"),
        ("💬", "Subject Messaging", "Direct communication between learners and teachers within subject contexts"),
        ("🤖", "AI Tutor", "Built-in AI assistant to guide students without giving direct answers"),
        ("🌙", "Dark Mode", "Full dark theme support for comfortable nighttime studying"),
    ]
    
    # Left column
    y_pos = Inches(1.8)
    for i, (icon, title, desc) in enumerate(features[:4]):
        # Feature box
        feat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y_pos, Inches(6), Inches(1.2))
        feat_box.fill.solid()
        feat_box.fill.fore_color.rgb = LIGHT_GRAY
        feat_box.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), y_pos + Inches(0.15), Inches(5.5), Inches(0.4))
        ttf = title_box.text_frame
        ttf.text = f"{icon} {title}"
        tp = ttf.paragraphs[0]
        tp.font.size = Pt(16)
        tp.font.bold = True
        tp.font.color.rgb = DARK_BG
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(0.8), y_pos + Inches(0.55), Inches(5.5), Inches(0.6))
        dtf = desc_box.text_frame
        dtf.word_wrap = True
        dtf.text = desc
        dp = dtf.paragraphs[0]
        dp.font.size = Pt(11)
        dp.font.color.rgb = RGBColor(100, 116, 139)
        
        y_pos += Inches(1.4)
    
    # Right column
    y_pos = Inches(1.8)
    for icon, title, desc in features[4:]:
        feat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), y_pos, Inches(6), Inches(1.2))
        feat_box.fill.solid()
        feat_box.fill.fore_color.rgb = LIGHT_GRAY
        feat_box.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(7.1), y_pos + Inches(0.15), Inches(5.5), Inches(0.4))
        ttf = title_box.text_frame
        ttf.text = f"{icon} {title}"
        tp = ttf.paragraphs[0]
        tp.font.size = Pt(16)
        tp.font.bold = True
        tp.font.color.rgb = DARK_BG
        
        desc_box = slide.shapes.add_textbox(Inches(7.1), y_pos + Inches(0.55), Inches(5.5), Inches(0.6))
        dtf = desc_box.text_frame
        dtf.word_wrap = True
        dtf.text = desc
        dp = dtf.paragraphs[0]
        dp.font.size = Pt(11)
        dp.font.color.rgb = RGBColor(100, 116, 139)
        
        y_pos += Inches(1.4)

def add_user_roles_slide():
    """User Roles Slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = ACCENT_ORANGE
    header.line.fill.background()
    
    header_tf = header.text_frame
    header_tf.text = "User Roles & Dashboards"
    hp = header_tf.paragraphs[0]
    hp.font.size = Pt(36)
    hp.font.bold = True
    hp.font.color.rgb = WHITE
    hp.alignment = PP_ALIGN.LEFT
    header_tf.margin_left = Inches(0.5)
    
    roles = [
        ("👨‍🎓", "Learners", PRIMARY_BLUE, [
            "Access learning materials & resources",
            "Take quizzes and submit assignments",
            "Track progress and view grades",
            "AI Tutor for personalized help",
            "Subject-specific messaging",
        ]),
        ("👨‍🏫", "Teachers", ACCENT_GREEN, [
            "Create and manage quizzes",
            "Upload learning materials",
            "Grade assignments & submissions",
            "Post announcements",
            "Track student progress",
        ]),
        ("👨‍💼", "Administrators", ACCENT_ORANGE, [
            "School & user management",
            "Approve teacher registrations",
            "Subject and grade configuration",
            "Global notifications",
            "Support ticket management",
        ]),
    ]
    
    x_pos = Inches(0.8)
    for icon, role, color, permissions in roles:
        # Role card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(2), Inches(3.9), Inches(5))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = color
        card.line.width = Pt(3)
        
        # Role header
        role_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(2), Inches(3.9), Inches(1))
        role_header.fill.solid()
        role_header.fill.fore_color.rgb = color
        role_header.line.fill.background()
        
        rh_tf = role_header.text_frame
        rh_tf.text = f"{icon} {role}"
        rhp = rh_tf.paragraphs[0]
        rhp.font.size = Pt(22)
        rhp.font.bold = True
        rhp.font.color.rgb = WHITE
        rhp.alignment = PP_ALIGN.CENTER
        
        # Permissions
        y_perm = Inches(3.3)
        for perm in permissions:
            perm_box = slide.shapes.add_textbox(x_pos + Inches(0.2), y_perm, Inches(3.5), Inches(0.5))
            ptf = perm_box.text_frame
            ptf.text = "✓ " + perm
            pp = ptf.paragraphs[0]
            pp.font.size = Pt(12)
            pp.font.color.rgb = DARK_BG
            y_perm += Inches(0.7)
        
        x_pos += Inches(4.2)

def add_architecture_slide():
    """System Architecture Slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(236, 72, 153)
    header.line.fill.background()
    
    header_tf = header.text_frame
    header_tf.text = "System Architecture"
    hp = header_tf.paragraphs[0]
    hp.font.size = Pt(36)
    hp.font.bold = True
    hp.font.color.rgb = WHITE
    hp.alignment = PP_ALIGN.LEFT
    header_tf.margin_left = Inches(0.5)
    
    # Architecture diagram boxes
    # Frontend
    frontend = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(1.8), Inches(3.3), Inches(1.2))
    frontend.fill.solid()
    frontend.fill.fore_color.rgb = PRIMARY_BLUE
    frontend.line.fill.background()
    
    ff_tf = frontend.text_frame
    ff_tf.text = "React Frontend\n(Vite + Tailwind)"
    ffp = ff_tf.paragraphs[0]
    ffp.font.size = Pt(16)
    ffp.font.bold = True
    ffp.font.color.rgb = WHITE
    ffp.alignment = PP_ALIGN.CENTER
    ff_tf.paragraphs[1].font.size = Pt(11)
    ff_tf.paragraphs[1].font.color.rgb = RGBColor(191, 219, 254)
    ff_tf.paragraphs[1].alignment = PP_ALIGN.CENTER
    
    # API
    api = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(3.3), Inches(3.3), Inches(0.8))
    api.fill.solid()
    api.fill.fore_color.rgb = ACCENT_ORANGE
    api.line.fill.background()
    
    api_tf = api.text_frame
    api_tf.text = "REST API (Axios)"
    api_p = api_tf.paragraphs[0]
    api_p.font.size = Pt(14)
    api_p.font.bold = True
    api_p.font.color.rgb = WHITE
    api_p.alignment = PP_ALIGN.CENTER
    
    # Backend
    backend = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(4.4), Inches(3.3), Inches(1.2))
    backend.fill.solid()
    backend.fill.fore_color.rgb = ACCENT_GREEN
    backend.line.fill.background()
    
    bb_tf = backend.text_frame
    bb_tf.text = "Node.js Backend\n(Express + Socket.IO)"
    bbp = bb_tf.paragraphs[0]
    bbp.font.size = Pt(16)
    bbp.font.bold = True
    bbp.font.color.rgb = WHITE
    bbp.alignment = PP_ALIGN.CENTER
    bb_tf.paragraphs[1].font.size = Pt(11)
    bb_tf.paragraphs[1].font.color.rgb = RGBColor(187, 247, 208)
    bb_tf.paragraphs[1].alignment = PP_ALIGN.CENTER
    
    # Database
    db = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(5.9), Inches(3.3), Inches(0.8))
    db.fill.solid()
    db.fill.fore_color.rgb = PRIMARY_PURPLE
    db.line.fill.background()
    
    db_tf = db.text_frame
    db_tf.text = "PostgreSQL Database"
    db_p = db_tf.paragraphs[0]
    db_p.font.size = Pt(14)
    db_p.font.bold = True
    db_p.font.color.rgb = WHITE
    db_p.alignment = PP_ALIGN.CENTER
    
    # External Services
    cloudinary = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9), Inches(2.5), Inches(3), Inches(0.8))
    cloudinary.fill.solid()
    cloudinary.fill.fore_color.rgb = RGBColor(59, 130, 246)
    cloudinary.line.fill.background()
    
    c_tf = cloudinary.text_frame
    c_tf.text = "Cloudinary\n(File Storage)"
    cp = c_tf.paragraphs[0]
    cp.font.size = Pt(12)
    cp.font.bold = True
    cp.font.color.rgb = WHITE
    cp.alignment = PP_ALIGN.CENTER
    c_tf.paragraphs[1].font.size = Pt(10)
    c_tf.paragraphs[1].font.color.rgb = RGBColor(191, 219, 254)
    c_tf.paragraphs[1].alignment = PP_ALIGN.CENTER
    
    # Features list on left
    features_title = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(0.5))
    ft_tf = features_title.text_frame
    ft_tf.text = "Key Technical Features:"
    ftp = ft_tf.paragraphs[0]
    ftp.font.size = Pt(16)
    ftp.font.bold = True
    ftp.font.color.rgb = DARK_BG
    
    tech_features = [
        "JWT Authentication",
        "Real-time WebSocket",
        "Responsive Design",
        "Dark Mode Support",
        "File Upload/Download",
        "Role-Based Routing",
        "API Interceptors",
        "Error Boundaries",
    ]
    
    y_feat = Inches(2.6)
    for feat in tech_features:
        feat_box = slide.shapes.add_textbox(Inches(0.5), y_feat, Inches(4), Inches(0.4))
        ftf = feat_box.text_frame
        ftf.text = "→ " + feat
        fp = ftf.paragraphs[0]
        fp.font.size = Pt(12)
        fp.font.color.rgb = RGBColor(71, 85, 105)
        y_feat += Inches(0.45)

def add_deployment_slide():
    """Deployment & Hosting Slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(6, 182, 212)
    header.line.fill.background()
    
    header_tf = header.text_frame
    header_tf.text = "Deployment & DevOps"
    hp = header_tf.paragraphs[0]
    hp.font.size = Pt(36)
    hp.font.bold = True
    hp.font.color.rgb = WHITE
    hp.alignment = PP_ALIGN.LEFT
    header_tf.margin_left = Inches(0.5)
    
    platforms = [
        ("Docker", "Containerized deployment with multi-stage builds", ACCENT_ORANGE),
        ("NGINX", "Production web server with optimized configuration", ACCENT_GREEN),
        ("Railway", "Backend hosting with automatic deployments", PRIMARY_PURPLE),
        ("Netlify/Vercel", "Frontend hosting with CDN distribution", PRIMARY_BLUE),
    ]
    
    y_pos = Inches(2)
    for platform, desc, color in platforms:
        # Platform box
        plat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(5.5), Inches(1))
        plat_box.fill.solid()
        plat_box.fill.fore_color.rgb = color
        plat_box.line.fill.background()
        
        plat_tf = plat_box.text_frame
        plat_tf.text = platform
        plat_p = plat_tf.paragraphs[0]
        plat_p.font.size = Pt(20)
        plat_p.font.bold = True
        plat_p.font.color.rgb = WHITE
        plat_tf.margin_left = Inches(0.3)
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(6.6), y_pos + Inches(0.2), Inches(6), Inches(0.6))
        dtf = desc_box.text_frame
        dtf.text = desc
        dp = dtf.paragraphs[0]
        dp.font.size = Pt(14)
        dp.font.color.rgb = RGBColor(71, 85, 105)
        
        y_pos += Inches(1.3)
    
    # Additional info
    info_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6), Inches(11.7), Inches(1))
    info_box.fill.solid()
    info_box.fill.fore_color.rgb = LIGHT_GRAY
    info_box.line.fill.background()
    
    info_tf = info_box.text_frame
    info_tf.text = "✓ Production-ready with environment-based configuration | ✓ Health checks enabled | ✓ Automated CI/CD ready"
    info_p = info_tf.paragraphs[0]
    info_p.font.size = Pt(14)
    info_p.font.color.rgb = RGBColor(71, 85, 105)
    info_p.alignment = PP_ALIGN.CENTER

def add_market_slide():
    """Market Opportunity Slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(139, 92, 246)
    header.line.fill.background()
    
    header_tf = header.text_frame
    header_tf.text = "Market Opportunity"
    hp = header_tf.paragraphs[0]
    hp.font.size = Pt(36)
    hp.font.bold = True
    hp.font.color.rgb = WHITE
    hp.alignment = PP_ALIGN.LEFT
    header_tf.margin_left = Inches(0.5)
    
    stats = [
        ("$350B+", "Global EdTech Market by 2025"),
        ("20%", "Annual Growth Rate (CAGR)"),
        ("1.6B+", "Students Worldwide"),
        ("85%", "Schools Seeking Digital Solutions"),
    ]
    
    x_pos = Inches(0.8)
    for stat, label in stats:
        # Stat card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(2), Inches(2.9), Inches(2))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.fill.background()
        
        # Stat number
        stat_box = slide.shapes.add_textbox(x_pos, Inches(2.3), Inches(2.9), Inches(0.8))
        stf = stat_box.text_frame
        stf.text = stat
        sp = stf.paragraphs[0]
        sp.font.size = Pt(36)
        sp.font.bold = True
        sp.font.color.rgb = PRIMARY_BLUE
        sp.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = slide.shapes.add_textbox(x_pos, Inches(3.2), Inches(2.9), Inches(0.6))
        ltf = label_box.text_frame
        ltf.text = label
        lp = ltf.paragraphs[0]
        lp.font.size = Pt(12)
        lp.font.color.rgb = RGBColor(71, 85, 105)
        lp.alignment = PP_ALIGN.CENTER
        
        x_pos += Inches(3.1)
    
    # Target markets
    markets_title = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.5))
    mt_tf = markets_title.text_frame
    mt_tf.text = "Target Markets:"
    mtp = mt_tf.paragraphs[0]
    mtp.font.size = Pt(20)
    mtp.font.bold = True
    mtp.font.color.rgb = DARK_BG
    
    markets = [
        ("🏫", "Primary & Secondary Schools", "Complete school management solution"),
        ("🎓", "Private Tutoring Centers", "Scalable multi-location support"),
        ("🏢", "Corporate Training", "Employee onboarding & skill development"),
        ("🌍", "Online Education Platforms", "White-label LMS capabilities"),
    ]
    
    y_market = Inches(5.1)
    for icon, market, desc in markets:
        market_box = slide.shapes.add_textbox(Inches(0.8), y_market, Inches(5.5), Inches(0.5))
        mktf = market_box.text_frame
        mktf.text = f"{icon} {market} - {desc}"
        mkp = mktf.paragraphs[0]
        mkp.font.size = Pt(14)
        mkp.font.color.rgb = RGBColor(71, 85, 105)
        y_market += Inches(0.6)

def add_benefits_slide():
    """Benefits Slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(20, 184, 166)
    header.line.fill.background()
    
    header_tf = header.text_frame
    header_tf.text = "Why Choose E-tab?"
    hp = header_tf.paragraphs[0]
    hp.font.size = Pt(36)
    hp.font.bold = True
    hp.font.color.rgb = WHITE
    hp.alignment = PP_ALIGN.LEFT
    header_tf.margin_left = Inches(0.5)
    
    benefits = [
        ("⚡", "Fast Performance", "Vite-powered build system for lightning-fast load times", PRIMARY_BLUE),
        ("🔒", "Secure by Design", "JWT authentication with role-based access control", ACCENT_GREEN),
        ("📱", "Mobile Ready", "Responsive design works on all devices and screen sizes", ACCENT_ORANGE),
        ("🤖", "AI Integration", "Built-in AI tutor enhances learning without replacing teachers", PRIMARY_PURPLE),
        ("💰", "Cost Effective", "Open-source technology stack reduces licensing costs", RGBColor(236, 72, 153)),
        ("🔧", "Easy to Maintain", "Modern React architecture with clear component structure", RGBColor(6, 182, 212)),
    ]
    
    positions = [
        (Inches(0.8), Inches(2)),
        (Inches(4.5), Inches(2)),
        (Inches(8.2), Inches(2)),
        (Inches(0.8), Inches(4.2)),
        (Inches(4.5), Inches(4.2)),
        (Inches(8.2), Inches(4.2)),
    ]
    
    for (x, y), (icon, title, desc, color) in zip(positions, benefits):
        # Benefit card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.5), Inches(2))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = color
        card.line.width = Pt(2)
        
        # Icon
        icon_box = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.2), Inches(0.6), Inches(0.6))
        icon_box.fill.solid()
        icon_box.fill.fore_color.rgb = color
        icon_box.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1), Inches(3.1), Inches(0.4))
        ttf = title_box.text_frame
        ttf.text = f"{icon} {title}"
        tp = ttf.paragraphs[0]
        tp.font.size = Pt(14)
        tp.font.bold = True
        tp.font.color.rgb = DARK_BG
        
        # Description
        desc_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.4), Inches(3.1), Inches(0.5))
        dtf = desc_box.text_frame
        dtf.word_wrap = True
        dtf.text = desc
        dp = dtf.paragraphs[0]
        dp.font.size = Pt(10)
        dp.font.color.rgb = RGBColor(100, 116, 139)

def add_roadmap_slide():
    """Roadmap Slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(99, 102, 241)
    header.line.fill.background()
    
    header_tf = header.text_frame
    header_tf.text = "Product Roadmap"
    hp = header_tf.paragraphs[0]
    hp.font.size = Pt(36)
    hp.font.bold = True
    hp.font.color.rgb = WHITE
    hp.alignment = PP_ALIGN.LEFT
    header_tf.margin_left = Inches(0.5)
    
    roadmap_items = [
        ("✅", "Current", "Core LMS features, quizzes, assignments, AI tutor", ACCENT_GREEN),
        ("🔄", "Q2 2024", "Mobile app (React Native), video conferencing", PRIMARY_BLUE),
        ("📅", "Q3 2024", "Advanced analytics, parent portal, payment integration", ACCENT_ORANGE),
        ("🔮", "2025", "VR/AR learning modules, blockchain certificates", PRIMARY_PURPLE),
    ]
    
    x_pos = Inches(0.8)
    for icon, phase, features, color in roadmap_items:
        # Phase card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(2.5), Inches(2.9), Inches(3.5))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = color
        card.line.width = Pt(3)
        
        # Phase header
        phase_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(2.5), Inches(2.9), Inches(0.8))
        phase_header.fill.solid()
        phase_header.fill.fore_color.rgb = color
        phase_header.line.fill.background()
        
        ph_tf = phase_header.text_frame
        ph_tf.text = f"{icon} {phase}"
        php = ph_tf.paragraphs[0]
        php.font.size = Pt(18)
        php.font.bold = True
        php.font.color.rgb = WHITE
        php.alignment = PP_ALIGN.CENTER
        
        # Features
        feat_box = slide.shapes.add_textbox(x_pos + Inches(0.2), Inches(3.5), Inches(2.5), Inches(2.3))
        ftf = feat_box.text_frame
        ftf.word_wrap = True
        ftf.text = features
        fp = ftf.paragraphs[0]
        fp.font.size = Pt(12)
        fp.font.color.rgb = RGBColor(71, 85, 105)
        
        # Arrow to next (except last)
        if x_pos < Inches(10):
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x_pos + Inches(2.95), Inches(4), Inches(0.4), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(203, 213, 225)
            arrow.line.fill.background()
        
        x_pos += Inches(3.1)

def add_contact_slide():
    """Contact/CTA Slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Gradient background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1))
    tf = title_box.text_frame
    tf.text = "Ready to Transform Education?"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11.333), Inches(0.8))
    stf = sub_box.text_frame
    stf.text = "Join schools worldwide using E-tab to deliver better learning experiences"
    sp = stf.paragraphs[0]
    sp.font.size = Pt(20)
    sp.font.color.rgb = RGBColor(148, 163, 184)
    sp.alignment = PP_ALIGN.CENTER
    
    # CTA Box
    cta_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(4.5), Inches(5.3), Inches(1.2))
    cta_box.fill.solid()
    cta_box.fill.fore_color.rgb = PRIMARY_BLUE
    cta_box.line.fill.background()
    
    cta_tf = cta_box.text_frame
    cta_tf.text = "Get Started Today"
    cta_p = cta_tf.paragraphs[0]
    cta_p.font.size = Pt(24)
    cta_p.font.bold = True
    cta_p.font.color.rgb = WHITE
    cta_p.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11.333), Inches(1))
    ctf = contact_box.text_frame
    ctf.text = "www.etab-lms.com | demo@etab-lms.com | +1 (555) 123-4567"
    cp = ctf.paragraphs[0]
    cp.font.size = Pt(16)
    cp.font.color.rgb = RGBColor(148, 163, 184)
    cp.alignment = PP_ALIGN.CENTER

# Generate all slides
add_title_slide()
add_overview_slide()
add_tech_stack_slide()
add_features_slide()
add_user_roles_slide()
add_architecture_slide()
add_deployment_slide()
add_market_slide()
add_benefits_slide()
add_roadmap_slide()
add_contact_slide()

# Save presentation
output_path = "E-tab_LMS_Marketing_Presentation.pptx"
prs.save(output_path)
print(f"[OK] Presentation saved successfully: {output_path}")
print(f"[INFO] Total slides: {len(prs.slides)}")
